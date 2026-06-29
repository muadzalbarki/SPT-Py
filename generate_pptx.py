
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0F, 0x17, 0x2A)
GOLD = RGBColor(0xD4, 0xAF, 0x37)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
DARK = RGBColor(0x02, 0x06, 0x17)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MID_GRAY = RGBColor(0x47, 0x55, 0x69)
DARK_GRAY = RGBColor(0x1E, 0x29, 0x3B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from pptx.oxml.ns import qn
        solidFill = shape.fill._fill
        srgb = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_elem = srgb.makeelement(qn('a:alpha'), {'val': str(alpha)})
            srgb.append(alpha_elem)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri', line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, is_bold, text_color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color if text_color else color
        p.font.bold = is_bold if is_bold else bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        from pptx.oxml.ns import qn
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
            pPr.append(lnSpc)
        spcPct = lnSpc.find(qn('a:spcPct'))
        if spcPct is None:
            spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
            lnSpc.append(spcPct)
    return txBox

def add_bullet_textbox(slide, left, top, width, height, items, font_size=15, color=WHITE, font_name='Calibri', line_spacing=1.3, bullet_char="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, is_bold = item
        else:
            text, is_bold = item, False
        p.text = f"{bullet_char} {text}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = is_bold
        p.font.name = font_name
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.find(qn('a:lnSpc'))
        if lnSpc is None:
            lnSpc = pPr.makeelement(qn('a:lnSpc'), {})
            pPr.append(lnSpc)
        spcPct = lnSpc.find(qn('a:spcPct'))
        if spcPct is None:
            spcPct = lnSpc.makeelement(qn('a:spcPct'), {'val': str(int(line_spacing * 100000))})
            lnSpc.append(spcPct)
    return txBox

def add_table(slide, left, top, width, height, headers, rows, header_color=NAVY, header_text_color=WHITE, body_color=DARK_GRAY, body_text_color=WHITE, alt_color=DARK):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = header_text_color
            paragraph.font.bold = True
            paragraph.font.name = 'Calibri'
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = alt_color if i % 2 == 1 else body_color
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = body_text_color
                paragraph.font.name = 'Calibri'
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape

def make_slide_title(slide, number, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), NAVY)
    num_str = f"{number:02d}"
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(0.8), Inches(0.9), num_str, font_size=28, color=GOLD, bold=True)
    add_textbox(slide, Inches(1.5), Inches(0.15), Inches(11), Inches(0.9), title, font_size=30, color=WHITE, bold=True)
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(0.05), GOLD)

def make_footer(slide):
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                "Sekretariat DPRD Kota Salatiga | SPT-Py v3.7",
                font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 1: JUDUL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(3.4), NAVY)
add_rect(slide, Inches(0), Inches(3.4), SLIDE_W, Inches(0.06), GOLD)

add_textbox(slide, Inches(1), Inches(0.6), Inches(11.3), Inches(1.5),
            "SISTEM OTOMATISASI SURAT DINAS",
            font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(1.0),
            "SPT-Py",
            font_size=36, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(2.6), Inches(11.3), Inches(0.7),
            "Sekretariat DPRD Kota Salatiga",
            font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.5),
            "Disusun oleh: Peserta Magang Prodi TI Fakultas Dakwah UIN Salatiga",
            font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.5),
            "Tahun 2026",
            font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: LATAR BELAKANG
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 2, "Latar Belakang")
make_footer(slide)

items = [
    ("Pembuatan Surat Tugas (SPT) masih dilakukan secara manual — staf menyalin template Word berulang kali, mengganti nama, tanggal, dan data peserta satu per satu", False),
    ("Data pegawai tersebar di file terpisah (Excel, Word, nota dinas) — tidak ada basis data terpusat untuk referensi cepat", False),
    ("Format surat sering tidak konsisten — perbedaan margin, jenis font, ukuran huruf, dan tata letak antar dokumen", False),
    ("Dokumentasi surat sulit dilacak — arsip fisik mudah hilang, pencarian surat lama memakan waktu", False),
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(4.8), items, font_size=16, line_spacing=1.6)

add_rect(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.5), DARK_GRAY)
add_textbox(slide, Inches(1), Inches(6.25), Inches(11.3), Inches(0.4),
            "DAMPAK: Proses lambat, rawan kesalahan manusia (human error), tidak efisien, dan menyulitkan audit",
            font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: SOLUSI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 3, "Solusi — SPT-Py")
make_footer(slide)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5),
            "SPT-Py adalah aplikasi desktop yang mengotomatiskan pembuatan Surat Tugas Perjalanan Dinas (SPT) DPRD.",
            font_size=18, color=WHITE)

items = [
    ("Wizard 4 langkah untuk generate surat — panduan interaktif dari pilih peserta hingga preview dan generate", True),
    ("CRUD data pegawai terpusat dengan database SQLite — satu sumber data untuk semua pegawai DPRD", False),
    ("Import & Export Excel — isi data pegawai massal dari file .xlsx atau backup ke Excel", False),
    ("Template .docx dengan placeholder { } — cukup buat template sekali, data diisi otomatis saat generate", False),
    ("Riwayat surat + Export PDF — semua surat tercatat, bisa diexport ke PDF untuk arsip digital", False),
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.3), Inches(11.7), Inches(4.2), items, font_size=15, line_spacing=1.5)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
            "Platform: Windows | Database: SQLite (lokal) | Mode: Dark Mode | Bahasa: Indonesia",
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: TARGET PENGGUNA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 4, "Target Pengguna")
make_footer(slide)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.5),
            "Sekretariat DPRD Kota Salatiga",
            font_size=22, color=WHITE, bold=True)

items = [
    ("Staf administrasi Sekretariat DPRD — pembuat surat tugas perjalanan dinas", False),
    ("4 Komisi dengan kebutuhan berbeda: Pimpinan, Komisi A, Komisi B, Komisi C", False),
    ("Masing-masing komisi: 5–7 anggota (total ~25 pegawai terdaftar)", False),
    ("Non-teknis: aplikasi dirancang dengan UI sederhana, wizard interaktif, tanpa perlu pengetahuan database atau pemrograman", False),
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.3), Inches(11.7), Inches(4.0), items, font_size=16, line_spacing=1.6)

# ============================================================
# SLIDE 5: FITUR — DATA PEGAWAI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 5, "Fitur — Data Pegawai")
make_footer(slide)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5),
            "Menu PEGAWAI — manajemen data pegawai DPRD",
            font_size=20, color=GOLD, bold=True)

items = [
    ("Tambah / Edit / Hapus pegawai — data: Nama, NIP, Jabatan, Pangkat/Gol, Komisi, No. HP", False),
    ("Dialog modal untuk tambah/edit — form terstruktur, validasi input (NIP unique, required field)", False),
    ("Pencarian real-time — filter berdasarkan nama, NIP, atau jabatan saat mengetik", False),
    ("Import Excel — upload file .xlsx, data pegawai langsung masuk database", False),
    ("Export Excel — backup data pegawai ke file .xlsx untuk arsip", False),
    ("Edit via double-click pada tabel — cepat, tanpa perlu cari tombol edit", False),
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.5), items, font_size=15, line_spacing=1.5)

# ============================================================
# SLIDE 6: FITUR — GENERATE SURAT (WIZARD)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 6, "Fitur — Generate Surat (Wizard 4 Langkah)")
make_footer(slide)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5),
            "Menu GENERATE — inti dari aplikasi",
            font_size=20, color=GOLD, bold=True)

steps = [
    ("STEP 1 — Pilih Peserta", "Checklist per komisi (Pimpinan, A, B, C) dengan collapsible section. Search filter nama/jabatan. Tombol 'Pilih Semua' per komisi."),
    ("STEP 2 — Surat Tugas", "Pilih Komisi, Nomor Surat (auto-generate), Tanggal Surat, Hari/Tanggal Kepergian, Dasar & Tujuan perjalanan dinas."),
    ("STEP 3 — Detail Kunjungan", "Lokasi (Kab/Kota, Provinsi), Tentang, Materi, Waktu kunjungan (Hari, Tanggal, Pukul), input Pendamping dinamis."),
    ("STEP 4 — Preview & Generate", "Ringkasan lengkap data surat. Klik 'Generate Surat' → progress bar → file .docx tersimpan + record di database."),
]

y_start = Inches(2.3)
for i, (title, desc) in enumerate(steps):
    y = y_start + Inches(i * 1.2)
    add_rect(slide, Inches(0.8), y, Inches(0.4), Inches(0.4), GOLD)
    add_textbox(slide, Inches(1.4), y - Inches(0.05), Inches(3.5), Inches(0.5), title, font_size=15, color=GOLD, bold=True)
    add_textbox(slide, Inches(5.0), y - Inches(0.05), Inches(7.5), Inches(1.0), desc, font_size=13, color=WHITE)
    if i < 3:
        add_textbox(slide, Inches(6.0), y + Inches(0.7), Inches(0.5), Inches(0.4), "↓", font_size=18, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: FITUR — RIWAYAT & EXPORT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 7, "Fitur — Riwayat Surat & Export PDF")
make_footer(slide)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5),
            "Menu RIWAYAT SURAT + EXPORT",
            font_size=20, color=GOLD, bold=True)

items = [
    ("Riwayat Surat — daftar semua surat yang telah digenerate (maks. 100 record terbaru)", False),
    ("Kolom informasi: No Surat, Tanggal, Template, Peserta, Status File", False),
    ("Sorting descending — surat terbaru tampil di paling atas", False),
    ("Tombol Refresh — reload data dari database tanpa restart aplikasi", False),
    ("Export PDF — pilih surat, klik Export, LibreOffice headless mengonversi .docx → .pdf", False),
    ("Dashboard — 2 kartu statistik: Total Pegawai & Riwayat Surat + tabel 10 surat terbaru", False),
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.3), Inches(11.7), Inches(3.5), items, font_size=15, line_spacing=1.5)

# ============================================================
# SLIDE 8: ALUR GENERATE SURAT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 8, "Alur Generate Surat")
make_footer(slide)

# Flowchart boxes
box_w = Inches(2.0)
box_h = Inches(1.0)
arrow_w = Inches(0.6)
gap = Inches(0.4)
total_w = box_w * 5 + arrow_w * 4 + gap * 8
start_x = (SLIDE_W - total_w) / 2
y_box = Inches(2.6)

steps_flow = ["Pilih\nPeserta", "Surat\nTugas", "Detail\nKunjungan", "Preview", "Generate\n.docx"]  # 5 steps

for i, step in enumerate(steps_flow):
    x = start_x + Inches(i * (2.0 + 0.6 + 0.4))
    shape = add_rect(slide, x, y_box, box_w, box_h, DARK_GRAY)
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(2)
    add_textbox(slide, x, y_box + Inches(0.1), box_w, box_h, step, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(steps_flow) - 1:
        ax = x + box_w + Inches(0.1)
        add_textbox(slide, ax, y_box + Inches(0.1), arrow_w, box_h, "→", font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# Info boxes below
info_items = [
    ("Input: Peserta\nterpilih", "Form: Komisi,\nNomor, Tanggal,\nDasar, Tujuan", "Form: Lokasi,\nWaktu, Materi,\nPendamping", "Ringkasan\ndata surat", "Output: file\n.docx + DB"),
]
for i, info in enumerate(info_items):
    pass

y_info = Inches(4.2)
info_data = [
    "Pilih pegawai\nper komisi",
    "Isi data\nsurat tugas",
    "Lokasi, waktu,\n& pendamping",
    "Verifikasi\nsebelum generate",
    "File .docx\n+ record DB",
]
for i, info in enumerate(info_data):
    x = start_x + Inches(i * (2.0 + 0.6 + 0.4))
    add_textbox(slide, x, y_info, box_w, Inches(0.8), info, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5), DARK_GRAY)
add_textbox(slide, Inches(1), Inches(5.65), Inches(11.3), Inches(0.4),
            "Keunggulan: Waktu pembuatan surat berkurang dari ±15 menit (manual) menjadi ±2 menit (otomatis)",
            font_size=14, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: TEKNOLOGI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 9, "Teknologi yang Digunakan")
make_footer(slide)

headers = ["Teknologi / Library", "Fungsi / Peran"]
rows = [
    ("Python 3.12+", "Bahasa pemrograman utama aplikasi"),
    ("PySide6 (Qt for Python)", "Framework GUI desktop — widget, layout, signal/slot"),
    ("SQLAlchemy ORM", "Penerjemah database — tulis Python, otomatis jadi SQL"),
    ("SQLite", "Database lokal — file .db, tanpa server, ringan"),
    ("python-docx + lxml", "Baca & manipulasi file .docx (XML) — ganti placeholder"),
    ("qtawesome", "Ikon Material Design — sidebar, tombol, header"),
    ("pandas + openpyxl", "Import & export data pegawai ke Excel (.xlsx)"),
    ("LibreOffice (headless)", "Konversi .docx → .pdf tanpa membuka jendela GUI"),
    ("XlsxWriter", "Menulis file Excel dengan format rapi"),
]
add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5), headers, rows,
          header_color=NAVY, body_color=DARK_GRAY, alt_color=DARK)

# ============================================================
# SLIDE 10: ARSITEKTUR
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 10, "Arsitektur Aplikasi")
make_footer(slide)

layers = [
    ("UI LAYER (Presentation)", "PySide6 Widgets", [
        "MainWindow — Sidebar + Topbar + QStackedWidget",
        "5 Pages: Dashboard, Pegawai, Generate, Riwayat, Settings",
        "Components: Table, Button, Card, Checklist, SearchBar",
    ], WHITE),
    ("BUSINESS LOGIC LAYER (Services)", "Python Services", [
        "DocumentService — orchestrator generate surat",
        "TemplateEngine — manipulasi .docx placeholder",
        "ParticipantFormatter — format daftar peserta",
        "PDFService — export PDF via LibreOffice",
        "ExcelService — import/export Excel",
    ], GOLD),
    ("DATA LAYER (Database)", "SQLAlchemy + SQLite", [
        "Models: Pegawai, Surat, TemplateDokumen, PesertaSurat",
        "Repository: CRUD operations (PegawaiRepo, SuratRepo)",
        "Session: SQLAlchemy session (auto-close context manager)",
        "Seeder: 25 pegawai awal + 1 template .docx",
    ], WHITE),
]

y_pos = Inches(1.8)
for i, (title, subtitle, items, accent_color) in enumerate(layers):
    layer_h = Inches(1.6)
    add_rect(slide, Inches(0.6), y_pos, Inches(12.1), layer_h, DARK_GRAY)
    add_rect(slide, Inches(0.6), y_pos, Inches(0.06), layer_h, accent_color)
    add_textbox(slide, Inches(0.9), y_pos + Inches(0.05), Inches(4), Inches(0.4), title, font_size=16, color=accent_color, bold=True)
    add_textbox(slide, Inches(0.9), y_pos + Inches(0.45), Inches(4), Inches(0.3), subtitle, font_size=11, color=LIGHT_GRAY)

    # Items in box
    item_text = "\n".join(f"• {item}" for item in items)
    add_textbox(slide, Inches(5.2), y_pos + Inches(0.05), Inches(7.2), Inches(1.4), item_text, font_size=11, color=WHITE)
    y_pos += layer_h + Inches(0.2)

# ============================================================
# SLIDE 11: STRUKTUR PROYEK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 11, "Struktur Proyek (Ringkas)")
make_footer(slide)

# Show the tree structure as formatted text
tree_text = (
    "SPT-Py/\n"
    "├── main.py                  # Entry point aplikasi\n"
    "├── app/                     # Kode inti aplikasi\n"
    "│   ├── app.py              # Startup orchestrator\n"
    "│   ├── config.py           # Konfigurasi path\n"
    "│   ├── database/           # Database layer\n"
    "│   │   ├── engine.py       # SQLAlchemy engine\n"
    "│   │   ├── models.py       # ORM models\n"
    "│   │   ├── repository.py   # CRUD operations\n"
    "│   │   └── seed.py         # Data awal\n"
    "│   ├── services/           # Business logic\n"
    "│   ├── themes/             # Dark mode tema\n"
    "│   ├── ui/                 # Antarmuka pengguna\n"
    "│   │   ├── main_window.py  # Kerangka jendela\n"
    "│   │   ├── components/     # Widget reusable\n"
    "│   │   ├── pages/          # 5 halaman aplikasi\n"
    "│   │   └── dialogs/        # Modal dialog\n"
    "│   └── utils/              # Fungsi bantu\n"
    "├── templates/               # Template .docx\n"
    "├── exports/                 # Output surat\n"
    "├── backups/                 # Riwayat versi\n"
    "├── requirements.txt          # Dependencies\n"
    "├── change.txt                # Riwayat perubahan\n"
    "├── alur.txt                  # Panduan pengguna\n"
    "├── dokumentasi-teknis.md     # Dokumentasi detail\n"
    "└── README.md                 # Dokumentasi publik"
)

add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.5), tree_text, font_size=11, color=WHITE, font_name='Consolas')

# ============================================================
# SLIDE 12: DEMO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 12, "Demo Aplikasi")
make_footer(slide)

add_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.5),
            "Screenshot Aplikasi",
            font_size=22, color=GOLD, bold=True)

# Placeholder boxes for screenshots
placeholder_data = [
    ("Wizard Step 1\nPilih Peserta", Inches(0.8), Inches(2.4), Inches(3.7), Inches(2.2)),
    ("Wizard Step 2\nSurat Tugas", Inches(4.8), Inches(2.4), Inches(3.7), Inches(2.2)),
    ("Wizard Step 3\nDetail Kunjungan", Inches(8.8), Inches(2.4), Inches(3.7), Inches(2.2)),
    ("Wizard Step 4\nPreview & Generate", Inches(0.8), Inches(4.9), Inches(3.7), Inches(2.2)),
    ("Tabel Pegawai\n+ Dialog", Inches(4.8), Inches(4.9), Inches(3.7), Inches(2.2)),
    ("Riwayat Surat\n+ Export PDF", Inches(8.8), Inches(4.9), Inches(3.7), Inches(2.2)),
]

for text, x, y, w, h in placeholder_data:
    shape = add_rect(slide, x, y, w, h, DARK_GRAY)
    shape.line.color.rgb = MID_GRAY
    shape.line.width = Pt(1)
    shape.line.dash_style = 2  # dash
    add_textbox(slide, x, y + Inches(0.6), w, Inches(1.0), text, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.3),
            "(Screenshot dapat ditambahkan dengan copy-paste gambar ke slide ini)",
            font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13: TANTANGAN & SOLUSI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 13, "Tantangan & Solusi")
make_footer(slide)

headers = ["Tantangan", "Solusi"]
rows = [
    ("Placeholder { } di .docx terpotong dalam beberapa XML run nodes, sehingga text replace biasa gagal", "TemplateEngine mengakses langsung XML dokumen, menggabungkan run nodes, baru melakukan replace"),
    ("Warna alternating row tabel secara default hijau (system accent) — tidak cocok dark mode navy", "Custom QSS: QTableWidget::item:alternate dengan token bg_table_alt"),
    ("Path template .docx berbeda antara Linux (development) dan Windows (produksi)", "Flexible path search di seed.py: cari di TEMPLATES_DIR, Downloads/, project root"),
    ("Checkbox checked tidak terlihat di dark mode (centang navy di background navy)", "Ubah QCheckBox::indicator:checked dari accent_navy ke accent_gold"),
    ("Light/dark toggle tidak diperlukan — memperumit kode, jarang digunakan", "Hapus toggle_theme(), set_light(), set_dark() — dark mode only. Kode lebih bersih, ~200 baris QSS berkurang"),
]
add_table(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.0), headers, rows,
          header_color=NAVY, body_color=DARK_GRAY, alt_color=DARK)

# ============================================================
# SLIDE 14: PENGEMBANGAN SELANJUTNYA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 14, "Pengembangan Selanjutnya")
make_footer(slide)

items = [
    ("Template Management — tambah/edit/hapus template surat langsung dari UI, tanpa perlu akses folder", False),
    ("Notifikasi — pemberitahuan otomatis ketika ada surat baru atau deadline mendekat", False),
    ("Multiple Template Types — dukung berbagai jenis surat (SPT, Surat Permohonan, Surat Kunjungan Kerja) dengan placeholder berbeda", False),
    ("Auto-numbering Cerdas — format nomor surat yang lebih fleksibel dan dapat dikonfigurasi pengguna", False),
    ("Export PDF Built-in — tanpa ketergantungan LibreOffice, menggunakan library Python murni", False),
]
add_bullet_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.0), items, font_size=16, line_spacing=1.6)

# ============================================================
# SLIDE 15: PENUTUP
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 15, "Penutup")
make_footer(slide)

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.5),
            "Kesimpulan",
            font_size=24, color=GOLD, bold=True)

conclusions = [
    ("SPT-Py berhasil mengotomatiskan pembuatan Surat Tugas Perjalanan Dinas di lingkungan Sekretariat DPRD Kota Salatiga", False),
    ("Proses pembuatan surat berkurang dari ±15 menit menjadi ±2 menit — efisiensi waktu hingga 85%", False),
    ("Data pegawai terpusat dalam satu database — mengurangi duplikasi dan inkonsistensi data", False),
    ("Format surat konsisten — template sekali, hasil seragam, margin dan font sesuai standar", False),
    ("Riwayat surat terdokumentasi dengan baik — memudahkan pencarian, audit, dan arsip digital", False),
]
add_bullet_textbox(slide, Inches(0.8), Inches(2.5), Inches(11.7), Inches(3.0), conclusions, font_size=15, line_spacing=1.5)

# Thank you
add_rect(slide, Inches(3.5), Inches(5.6), Inches(6.3), Inches(0.8), NAVY)
add_textbox(slide, Inches(3.5), Inches(5.7), Inches(6.3), Inches(0.6),
            "TERIMA KASIH",
            font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4),
            "Semoga aplikasi ini bermanfaat bagi kelancaran administrasi Sekretariat DPRD Kota Salatiga",
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 16: Q&A / KONTAK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
make_slide_title(slide, 16, "Q&A / Sesi Tanya Jawab")
make_footer(slide)

add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.6),
            "Sesi Tanya Jawab",
            font_size=28, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(3.5), Inches(2.8), Inches(6.3), Inches(2.5), DARK_GRAY)
add_textbox(slide, Inches(3.8), Inches(2.9), Inches(5.7), Inches(0.4),
            "Kontak Penyusun",
            font_size=18, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

contact_lines = [
    ("Nama", "Peserta Magang Prodi TI", True),
    ("Program Studi", "Teknik Informatika", True),
    ("Fakultas", "Dakwah", True),
    ("Universitas", "UIN Salatiga", True),
    ("Tahun", "2026", True),
]
y = Inches(3.4)
for label, value, _ in contact_lines:
    add_textbox(slide, Inches(4.0), y, Inches(2.5), Inches(0.35), label, font_size=13, color=MID_GRAY)
    add_textbox(slide, Inches(6.5), y, Inches(5), Inches(0.35), f":  {value}", font_size=13, color=WHITE)
    y += Inches(0.35)

add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
            "Terima kasih atas perhatian dan partisipasi Bapak/Ibu",
            font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SAVE
# ============================================================
output_path = r"C:\Users\lenov\Documents\SPT-Py_Presentasi.pptx"
prs.save(output_path)
print(f"Presentasi berhasil dibuat: {output_path}")
